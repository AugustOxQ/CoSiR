#!/usr/bin/env python3
"""Build the 2026-08-26 Conditional Buddies weekly report slide deck."""
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

OUT = Path(__file__).resolve().parents[1] / "2026-08-26_weekly_conditional_buddies_slides.pptx"
W, H = Inches(13.333), Inches(7.5)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x59, 0x59, 0x59)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
BORDER = RGBColor(0xBF, 0xBF, 0xBF)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
RED = RGBColor(0xC0, 0x30, 0x30)
WHITE = RGBColor(255, 255, 255)


def set_font(paragraph, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def text(slide, value, x, y, w, h, size=12, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    set_font(p, size, bold, color, align)
    return box


def bullet_text(slide, items, x, y, w, h, size=14, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(9)
        set_font(p, size, False, color)
        p.text = "•  " + item
    return box


def border_cell(cell):
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tc_pr.append(parse_xml(
            '<%s %s w="12700" cap="flat" cmpd="sng" algn="ctr">'
            '<a:solidFill><a:srgbClr val="BFBFBF"/></a:solidFill>'
            '<a:prstDash val="solid"/><a:round/><a:headEnd type="none" w="med" len="med"/>'
            '<a:tailEnd type="none" w="med" len="med"/></%s>' % (edge, nsdecls('a'), edge)))


def table(slide, headers, rows, x, y, w, h, widths=None, font_size=12):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    t = shape.table
    if widths:
        for col, width in zip(t.columns, widths): col.width = Inches(width)
    for r, values in enumerate([headers] + rows):
        for c, value in enumerate(values):
            cell = t.cell(r, c)
            cell.text = str(value)
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r == 0 else WHITE
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            numeric = c > 0
            set_font(p, font_size, r == 0, BLACK, PP_ALIGN.CENTER if numeric else PP_ALIGN.LEFT)
            border_cell(cell)
    return shape


def bar_chart(slide, categories, series, x, y, w, h, title=None, minimum=None, maximum=None):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Calibri"; chart.legend.font.size = Pt(10)
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        set_font(chart.chart_title.text_frame.paragraphs[0], 12, True)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.name = "Calibri"; chart.value_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = "Calibri"; chart.category_axis.tick_labels.font.size = Pt(10)
    if minimum is not None: chart.value_axis.minimum_scale = minimum
    if maximum is not None: chart.value_axis.maximum_scale = maximum
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.plots[0].data_labels.font.name = "Calibri"; chart.plots[0].data_labels.font.size = Pt(9)
    colors = [BLUE, ORANGE, RED]
    for s, color in zip(chart.series, colors):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = color
        s.format.line.color.rgb = color
    return chart


def base_slide(prs, kicker, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    text(slide, kicker.upper(), 0.5, 0.28, 10.7, 0.25, 11, False, GRAY)
    text(slide, title_text, 0.5, 0.55, 12.25, 0.55, 24, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.18), Inches(12.33), Inches(0.012))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER; line.line.fill.background()
    return slide


def footer(slide, page, total):
    text(slide, f"{page} / {total}", 12.2, 7.15, 0.9, 0.2, 10, False, GRAY, PP_ALIGN.RIGHT)


def add_note(slide, value, y=5.85, h=0.8):
    return text(slide, value, 0.5, y, 12.25, h, 12, False, BLACK)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    total = 16

    # 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    text(s, "Conditional Buddies — Weekly Results", 0.7, 2.6, 11.9, 0.65, 32, True)
    text(s, "What matters after initialization?", 0.7, 3.35, 11.9, 0.4, 18, False, GRAY)
    text(s, "RedCaps-150k  ·  graph source  ·  edge typing  ·  condition-table behavior", 0.7, 4.9, 11.9, 0.3, 14)
    text(s, "2026-08-26  ·  branch experiment/condition_drift_retrieval_correlation", 0.7, 5.3, 11.9, 0.3, 12, False, GRAY)
    footer(s, 1, total)

    # 2
    s = base_slide(prs, "This week", "From “buddy-init works” to “which pieces matter?”")
    text(s, "Last week", 0.5, 1.5, 2.2, 0.25, 13, True)
    text(s, "Buddy-graph initialization beat the generic initializer on RedCaps, including at 300k — but neither conditioned approach beat raw CLIP retrieval.", 0.5, 1.8, 5.7, 0.9, 15)
    text(s, "This week", 6.7, 1.5, 2.2, 0.25, 13, True)
    text(s, "Ablated graph source and edge weighting; explained RedCaps lift scaling; separated held-out retrieval from assigned training-row behavior.", 6.7, 1.8, 5.6, 0.9, 15)
    table(s, ["New result", "Bottom line"], [
        ["Experiment 8", "Source pair matters for i2t, not t2i"],
        ["Experiment 9", "Signal is stable; apparent size effect is normalization"],
        ["Experiment 10", "Typed-edge correction changes init, not retrieval"],
        ["Experiments 11.1–11.2", "Frozen wins held-out; trained helps its in-sample task"],
    ], 0.5, 3.15, 12.25, 1.9, [2.45, 9.8], 12)
    add_note(s, "All retrieval reads use paired-within-seed deltas and mean/SEM; measured noise floor: ~0.1–0.7 R1.", 5.55, 0.45)
    footer(s, 2, total)

    # 3
    s = base_slide(prs, "Experiment 8 · source-pair ablation", "The graph’s source pair matters — selectively")
    text(s, "Fixed frozen CLIP backbone/settings; all 16 cached vision × text encoder pairs; 3 seeds each (48 runs). Deltas vs. clip_img:clip_txt.", 0.5, 1.4, 12.25, 0.4, 12, False, GRAY)
    table(s, ["Encoder pair", "Mean Δ i2t R1", "mean/SEM", "Wins/3"], [
        ["dinov2:minilm", "+2.30", "+6.5", "3"], ["dinov2:bge", "+1.97", "+4.5", "3"],
        ["vit_sup:clip_txt", "+1.77", "+4.0", "3"], ["clip_img:bge", "+1.43", "+6.6", "3"],
        ["siglip_v:e5", "−0.37", "−0.8", "1"], ["vit_sup:e5", "−0.80", "−6.9", "0"],
    ], 0.5, 1.95, 5.65, 3.35, [2.4, 1.25, 1.1, 0.9], 12)
    bar_chart(s, ["dino / minilm", "dino / bge", "vit / clip", "clip / bge", "siglip / e5", "vit / e5"],
              [("Mean Δ i2t R1", [2.30, 1.97, 1.77, 1.43, -0.37, -0.80])], 6.55, 1.85, 6.2, 3.9, "i2t source-pair delta", -1.2, 2.7)
    footer(s, 3, total)

    # 4
    s = base_slide(prs, "Experiment 8 · interpretation", "i2t is sensitive; t2i is essentially invariant")
    bullet_text(s, [
        "i2t: 12/15 non-baseline pairs are positive and significant; all have 3/3 wins, and 10 are above the noise floor.",
        "t2i: only dinov2:minilm clears |mean/SEM| ≥ 2 (−0.07 R1; mean/SEM −2.0), still below the ~0.1–0.7 R1 noise floor.",
        "Survival is not a usefulness proxy: Pearson r = +0.010 (t2i) and −0.410 (i2t), n = 15.",
    ], 0.7, 1.65, 11.6, 2.3, 16)
    text(s, "Verdict", 0.7, 4.65, 2.0, 0.25, 13, True)
    text(s, "CLIP is not the strongest i2t source in this 150k grid. The best graph-source choice can matter, but it does not move t2i meaningfully.", 0.7, 4.98, 11.5, 0.65, 16, True)
    footer(s, 4, total)

    # 5
    s = base_slide(prs, "Experiment 9 · stability", "The subreddit signal is stable across scale")
    text(s, "Same-subreddit buddy-edge lift, independently sampled across all 350 subreddits.", 0.5, 1.42, 12.25, 0.3, 13, False, GRAY)
    table(s, ["Metric", "150k", "300k", "500k"], [
        ["Overall lift", "22.80×", "22.74×", "22.79×"],
        ["Lift-qualifying subreddits", "159 (45%)", "197 (56%)", "214 (61%)"],
        ["z-qualifying subreddits", "124 (35%)", "157 (45%)", "185 (53%)"],
        ["Lift–size Spearman ρ", "−0.523", "−0.532", "−0.590"],
        ["z–size Spearman ρ", "+0.221", "+0.163", "+0.226"],
    ], 0.5, 1.95, 7.0, 3.2, [2.8, 1.4, 1.4, 1.4], 12)
    bar_chart(s, ["150k", "300k", "500k"], [("Overall lift (×)", [22.80, 22.74, 22.79])], 8.0, 2.0, 4.7, 3.15, "Aggregate lift", 0, 25)
    add_note(s, "At 150k, lift is 22.80× across 159/350 qualifying subreddits; individual lift spans 4.41× (pics) to 670.79× (f1porn).", 5.55, 0.55)
    footer(s, 5, total)

    # 6
    s = base_slide(prs, "Experiment 9 · explanation", "The apparent size story is mechanical, not topical dilution")
    table(s, ["Property vs. lift, 150k", "Pearson r", "Spearman ρ", "Read"], [
        ["Size", "−0.328", "−0.523", "real curved/rank relationship"],
        ["Caption diversity", "−0.219", "−0.224", "null"],
        ["Visual homogeneity", "+0.204", "+0.230", "null"],
    ], 0.5, 1.55, 12.25, 1.8, [3.4, 1.6, 1.7, 5.55], 12)
    bullet_text(s, [
        "Purity rises slightly with size; the data do not show that large subreddits are topically diluted.",
        "Lift falls because its denominator is degree-dependent: Spearman(deg_s, size) = +0.973.",
        "Positive z–size is the expected confidence effect, not evidence of a stronger underlying signal.",
    ], 0.7, 3.85, 11.5, 1.65, 16)
    text(s, "Verdict: C1’s aggregate signal is stable.", 0.7, 5.8, 11.5, 0.35, 16, True)
    footer(s, 6, total)

    # 7
    s = base_slide(prs, "Experiment 10 · edge typing", "Fixing edge typing moves the initializer, not retrieval")
    text(s, "Typed uses the supporting modality rank for image-only/text-only union-graph edges (~98% of edges); prior blend uses both modalities for every edge. 2 modes × 3 seeds.", 0.5, 1.4, 12.25, 0.45, 12, False, GRAY)
    table(s, ["Paired metric, typed − blend", "Mean Δ", "± std", "mean/SEM", "Typed wins"], [
        ["test_oracle/t2i_R1", "+0.00", "0.26", "+0.0", "2/3"],
        ["test_oracle/i2t_R1", "−0.37", "0.59", "−1.1", "1/3"],
        ["test_pre_diff/t2i_R1", "+0.00", "0.10", "+0.0", "1/3"],
        ["test_pre_diff/i2t_R1", "+0.17", "0.55", "+0.5", "2/3"],
    ], 0.5, 2.0, 7.15, 2.55, [3.1, 1, 0.9, 1.05, 1.1], 12)
    bar_chart(s, ["oracle t2i", "oracle i2t", "pre-diff t2i", "pre-diff i2t"], [("Mean Δ R1", [0.00, -0.37, 0.00, 0.17])], 8.1, 2.0, 4.65, 2.9, "Retrieval deltas", -0.6, 0.4)
    add_note(s, "test_raw was identical in all six runs: t2i_R1 = 28.1; i2t_R1 = 29.7.", 5.15, 0.35)
    footer(s, 7, total)

    # 8
    s = base_slide(prs, "Experiment 10 · interpretation", "The correction is real — the retrieval result is a citable null")
    bullet_text(s, [
        "Mean |typed − blend| = 0.5117: 1.02× the mean absolute embedding value (0.5000).",
        "Per-dimension correlation ranges from +0.985 to −0.996, with mean +0.240.",
        "Despite this substantial initialization change, retrieval and the gap to CLIP are unchanged within significance/noise criteria at 150k.",
    ], 0.7, 1.65, 11.5, 2.15, 16)
    text(s, "Verdict", 0.7, 4.65, 2.0, 0.25, 13, True)
    text(s, "Typed edges fix a genuine graph-level flaw, but buddy-derived structure appears more load-bearing than this fine edge-weighting choice.", 0.7, 4.98, 11.5, 0.7, 16, True)
    footer(s, 8, total)

    # 9
    s = base_slide(prs, "Experiment 11.1 · held-out retrieval", "Freeze the table to win held-out oracle i2t")
    text(s, "Post-init condition table: trained (em_interval = −1) vs. frozen (em_interval = 101, beyond 100 epochs). Buddy init, backbone, and loss stack fixed; 2 arms × 3 seeds.", 0.5, 1.4, 12.25, 0.45, 12, False, GRAY)
    table(s, ["Metric, frozen − trained", "Mean Δ", "mean/SEM", "Per-seed Δ", "Read"], [
        ["test_oracle/t2i_R1", "−0.27", "−2.0", "−0.40 / −0.40 / +0.00", "inside noise floor"],
        ["test_oracle/i2t_R1", "+4.67", "+32.1", "+4.90 / +4.40 / +4.70", "decisive frozen win"],
    ], 0.5, 2.0, 12.25, 1.55, [3.05, 1.2, 1.25, 3.75, 3.0], 12)
    bar_chart(s, ["oracle t2i", "oracle i2t"], [("Frozen − trained R1", [-0.27, 4.67])], 3.0, 4.05, 7.2, 2.0, "Held-out oracle delta", -1, 5.5)
    footer(s, 9, total)

    # 10
    s = base_slide(prs, "Experiment 11.1 · mechanism", "Training changes the table — and still transfers worse")
    table(s, ["Final-epoch geometry, trained − frozen", "Mean Δ", "mean/SEM"], [
        ["shift_mean", "+0.0081", "+13.9"], ["shift_std", "+0.0079", "+31.9"],
        ["Row diversity", "+0.0141", "+15.9"], ["Column diversity", "−0.0127", "−10.2"],
        ["Conditioned effective dimensions", "0", "exactly 301 in all six runs"],
    ], 0.5, 1.55, 8.2, 3.05, [4.5, 1.45, 2.25], 12)
    text(s, "Sanity check", 9.15, 1.65, 2.3, 0.25, 13, True)
    text(s, "Frozen drift from init: exactly 0 for all seeds.\n\nTrained drift: mean 0.0845; range [0.0834, 0.0855].", 9.15, 2.0, 3.2, 1.4, 15)
    add_note(s, "Verdict: training is not inert — it differentiates conditions and shifts embeddings further — but freezing decisively wins held-out oracle i2t; t2i is a noise-floor null.", 5.25, 0.7)
    footer(s, 10, total)

    # 11
    s = base_slide(prs, "Experiment 11.2 · isolated drift", "Drift has a large, graded rank effect when isolated")
    text(s, "For each checkpoint seed: 3,000 in-sample queries, own assigned condition, full 150,000-item training gallery. Counterfactual holds trained combiner, other_proj, and gallery fixed; replaces only trained conditions with buddy init.", 0.5, 1.38, 12.25, 0.5, 12, False, GRAY)
    table(s, ["Spearman ρ; n = 3000 / seed", "Seed 1", "Seed 2", "Seed 3"], [
        ["rho(delta_rank, drift)", "+0.026", "+0.017", "+0.009"],
        ["rho(|delta_rank|, drift)", "+0.160", "+0.141", "+0.156"],
        ["rho(|delta_rank|, shift)", "−0.288", "−0.296", "−0.302"],
        ["rho(delta_rank_swap, drift)", "+0.466", "+0.466", "+0.477"],
    ], 0.5, 2.05, 7.25, 2.55, [3.6, 1.2, 1.2, 1.2], 12)
    bar_chart(s, ["Seed 1", "Seed 2", "Seed 3"], [("rho(delta_rank_swap, drift)", [0.466, 0.466, 0.477])], 8.15, 2.05, 4.6, 2.55, "Condition-only effect", 0, 0.55)
    add_note(s, "Every condition-only rho has p < 1e−160. Mean delta_rank_swap: +343.5 / +352.0 / +357.1; medians: +6 / +7 / +7.", 5.2, 0.55)
    footer(s, 11, total)

    # 12
    s = base_slide(prs, "Experiment 11.2 · assigned training rows", "Trained wins its own in-sample, own-condition task")
    table(s, ["Seed", "Mean rank frozen / trained", "R1 frozen / trained", "Fraction improved"], [
        ["1", "1118.1 / 975.0", "0.0923 / 0.1007", "0.636"],
        ["2", "1122.8 / 980.4", "0.0880 / 0.0997", "0.640"],
        ["3", "1112.5 / 996.2", "0.0893 / 0.1033", "0.639"],
    ], 0.5, 1.65, 12.25, 2.05, [1.1, 4.0, 3.7, 3.45], 13)
    bar_chart(s, ["Seed 1", "Seed 2", "Seed 3"], [("Frozen R1", [0.0923, 0.0880, 0.0893]), ("Trained R1", [0.1007, 0.0997, 0.1033])], 3.0, 4.2, 7.0, 1.85, "Assigned-row R1", 0.08, 0.11)
    add_note(s, "Scope boundary: this is neither a held-out result nor oracle-max-over-conditions retrieval. It shows trained conditions help their assigned training rows — not that the learned table generalizes better as a held-out codebook.", 5.95, 0.55)
    footer(s, 12, total)

    # 13
    s = base_slide(prs, "Experiments 11.1–11.2 · synthesis", "The apparent conflict disappears when the constructs stay separate")
    table(s, ["Experiment 11.1", "Experiment 11.2"], [
        ["Held-out queries", "In-sample training rows"],
        ["Oracle max over conditions", "Own assigned condition"],
        ["Frozen beats trained on i2t: +4.67 R1", "Trained benefit graded by drift: ρ = +0.466 / +0.466 / +0.477"],
    ], 0.5, 1.65, 12.25, 2.0, [6.12, 6.13], 14)
    text(s, "Simple reading", 0.7, 4.3, 2.0, 0.25, 13, True)
    text(s, "Training learns useful sample-specific co-adaptation for known rows, but its resulting table/codebook transfers worse to unseen queries when those queries may select their best table entry.", 0.7, 4.62, 11.6, 0.7, 16)
    add_note(s, "These are not competing estimates of one quantity: the held-out-vs-in-sample and oracle-vs-own-condition mismatches mean 11.2 does not explain away 11.1’s frozen-wins result.", 5.7, 0.55)
    footer(s, 13, total)

    # 14
    s = base_slide(prs, "Experiment 11.3 · pending", "The coupling test is implemented; results are pending")
    text(s, "Question: should predictor-distillation remain one-way, or also pull the trainable table toward what condition_predictor can represent?", 0.5, 1.45, 12.25, 0.45, 15)
    bullet_text(s, [
        "loss.pred_stopgrad=True (default): table detached; gradients flow table → predictor.",
        "loss.pred_stopgrad=false: predictor loss also updates the trainable table.",
        "pred_coupled will run three RedCaps-150k seeds against 11.1’s trained/frozen arms, sharing buddy-init template and wandb group.",
    ], 0.7, 2.15, 11.5, 1.7, 15)
    table(s, ["Implementation boundary", "Evidence"], [
        ["Toggle", "038f956"], ["Sweep runner", "be0efb7 — one file, 85 insertions"],
        ["Paired analysis", "c821513"], ["Results commit", "None in requested git log --oneline -10"],
    ], 0.5, 4.25, 12.25, 1.6, [3.2, 9.05], 12)
    footer(s, 14, total)

    # 15
    s = base_slide(prs, "Next", "Close the coupling test, then target the remaining gates")
    bullet_text(s, [
        "Complete Experiment 11.3: report paired pred_coupled − trained/frozen deltas for test_oracle and test_pre_diff, both directions, plus drift and predictor loss.",
        "Test Experiment 8’s strongest i2t source pairs on test_pre_diff; extend to 300k only if held-out feature re-extraction is justified.",
        "Extend Experiment 9 to 1M/3.1M; fit a joint model only if explanation, not coverage, is the goal.",
        "Treat Experiment 10 as a completed 150k null unless a scale or operating-point rationale emerges.",
        "Keep Experiment 11 claims separate by construct; a held-out-reachable predictor diagnostic is needed to bridge the gap.",
    ], 0.7, 1.55, 11.55, 4.8, 14)
    footer(s, 15, total)

    # 16
    s = base_slide(prs, "Questions", "The publication-safe claim and the sharper open question")
    text(s, "Current publication-safe claim", 0.7, 1.65, 4.2, 0.25, 13, True)
    text(s, "Buddy-graph structure is a robust, content-grounded initializer and a better in-model starting point than the generic alternative.", 0.7, 2.0, 11.4, 0.65, 18, True)
    text(s, "This week’s sharper question", 0.7, 3.75, 4.2, 0.25, 13, True)
    text(s, "How do we retain useful trained-row adaptation without sacrificing held-out codebook behavior?", 0.7, 4.1, 11.4, 0.65, 20, True, BLUE)
    text(s, "No outcome should be inferred from the pending coupling implementation alone.", 0.7, 5.75, 11.4, 0.3, 13, False, GRAY)
    footer(s, 16, total)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
